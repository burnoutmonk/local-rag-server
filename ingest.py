from __future__ import annotations

import hashlib
import json
import re
import time
import uuid
from pathlib import Path
from typing import List, Tuple

from tqdm import tqdm
from pypdf import PdfReader
import docx

import torch
import os
from concurrent.futures import ProcessPoolExecutor, as_completed

from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams, PointStruct, Filter, FieldCondition, MatchValue

from config import (
    QDRANT_HOST, QDRANT_PORT, COLLECTION,
    EMBED_MODEL_NAME,
    MAX_CHARS, OVERLAP_CHARS, BATCH_SIZE,
)

DATA_DIR = Path(__file__).resolve().parent / "data_raw"
# Store hash file in /app/state if it exists (Docker volume), otherwise next to ingest.py
_STATE_DIR = Path("/app/state") if Path("/app/state").exists() else Path(__file__).resolve().parent
HASH_FILE = _STATE_DIR / ".ingest_hashes.json"


# ── File hashing ──────────────────────────────────────────────────────────────
def file_hash(path: Path) -> str:
    """MD5 hash of file contents — fast enough for large PDFs."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def load_hashes() -> dict:
    try:
        if HASH_FILE.exists() and HASH_FILE.is_file():
            text = HASH_FILE.read_text().strip()
            if text:
                return json.loads(text)
    except Exception:
        pass
    return {}


def save_hashes(hashes: dict) -> None:
    HASH_FILE.write_text(json.dumps(hashes, indent=2))


# ── Text processing ───────────────────────────────────────────────────────────
def normalize_text(s: str) -> str:
    s = s.replace("\u00a0", " ")
    s = s.replace("\f", "\n")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def clean_pdf_text(text: str) -> str:
    lines = text.splitlines()
    cleaned = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            cleaned.append("")
            continue
        if re.fullmatch(r"[-\u2013]?\s*\d+\s*[-\u2013]?", stripped):
            continue
        if re.fullmatch(r"[Pp]age\s+\d+(\s+of\s+\d+)?", stripped):
            continue
        if re.fullmatch(r"\d+\s*/\s*\d+", stripped):
            continue
        if re.search(r"\xa9|copyright|\bconfidential\b|\ball rights reserved\b|\bproprietary\b", stripped, re.IGNORECASE):
            continue
        if re.search(r"\.{4,}\s*\d+\s*$", stripped):
            continue
        if re.fullmatch(r"[.\-_\s]{5,}", stripped):
            continue
        cleaned.append(line)

    result = "\n".join(cleaned)
    result = re.sub(r"(\w)-\n(\w)", r"\1\2", result)
    result = re.sub(r"\n{3,}", "\n\n", result)
    return result.strip()


def chunk_text(text: str, max_chars: int = MAX_CHARS, overlap: int = OVERLAP_CHARS) -> List[str]:
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    chunks: List[str] = []
    buf = ""

    def flush() -> None:
        nonlocal buf
        if buf.strip():
            chunks.append(buf.strip())
        buf = ""

    def hard_split(prefix: str, long_para: str) -> None:
        start = 0
        while start < len(long_para):
            head = (prefix + "\n\n") if (prefix and start == 0) else ""
            available = max_chars - len(head)
            end = min(start + available, len(long_para))
            chunks.append((head + long_para[start:end]).strip())
            next_start = end - overlap
            if next_start <= start:
                next_start = end
            start = next_start

    for p in paras:
        if len(buf) + len(p) + 2 <= max_chars:
            buf = (buf + "\n\n" + p).strip() if buf else p
        else:
            overlap_prefix = buf[-overlap:] if (buf and overlap > 0) else ""
            flush()
            if len(p) <= max_chars:
                buf = (overlap_prefix + "\n\n" + p).strip() if overlap_prefix else p
            else:
                hard_split(overlap_prefix, p)

    flush()
    return chunks


def read_pdf_sections(path: Path) -> List[Tuple[str, str]]:
    try:
        reader = PdfReader(str(path), strict=False)
    except Exception as exc:
        print(f"    WARNING: could not open PDF ({exc}) -- skipping entire file")
        return []

    total_pages = len(reader.pages)
    out: List[Tuple[str, str]] = []
    for i, page in enumerate(reader.pages):
        # per-page prints will be interleaved under multiprocessing but still usable
        print(f"    [{path.name}] reading page {i+1}/{total_pages} ...", end="\r", flush=True)
        try:
            raw = normalize_text(page.extract_text() or "")
            text = clean_pdf_text(raw)
        except Exception as exc:
            print(f"    WARNING: page {i+1} of {path.name} failed ({exc}) -- skipping")
            continue
        if text:
            out.append((f"Page {i+1}", text))
    print(f"    [{path.name}] -> {len(out)}/{total_pages} pages with text    ")
    return out


def read_docx_sections(path: Path) -> List[Tuple[str, str]]:
    d = docx.Document(str(path))
    sections: List[Tuple[str, List[str]]] = []
    current_title = "Document"
    current_lines: List[str] = []

    for p in d.paragraphs:
        style = (p.style.name or "").lower()
        try:
            txt = normalize_text(p.text)
        except Exception as exc:
            print(f"    WARNING: paragraph parse failed in {path.name} ({exc}) -- skipping")
            continue
        if not txt:
            continue
        if "heading" in style:
            if current_lines:
                sections.append((current_title, current_lines))
                current_lines = []
            current_title = txt
        else:
            current_lines.append(txt)

    if current_lines:
        sections.append((current_title, current_lines))

    result = [(title, "\n".join(lines)) for title, lines in sections]
    print(f"    [{path.name}] -> {len(result)} sections found")
    return result


def read_txt_sections(path: Path) -> List[Tuple[str, str]]:
    """Read plain text or markdown file as a single section."""
    try:
        text = path.read_text(encoding="utf-8")
        text = normalize_text(text)
        if text:
            return [("Document", text)]
        return []
    except Exception as exc:
        print(f"    WARNING: could not read {path.name} ({exc}) -- skipping")
        return []


def read_csv_sections(path: Path) -> List[Tuple[str, str]]:
    """Read CSV file and format rows as readable text."""
    import csv
    try:
        with open(path, "r", encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)

        if not rows:
            return []

        # Check if we have a header
        header = rows[0] if len(rows) > 0 else []
        data_rows = rows[1:] if len(rows) > 1 else []

        lines = []
        for row_idx, row in enumerate(data_rows, 1):
            if len(header) > 20:
                # Wide CSV: use key=value format
                row_text = f"Row {row_idx}: " + ", ".join(
                    f"{header[i]}={row[i] if i < len(row) else ''}"
                    for i in range(len(header))
                )
            else:
                # Normal CSV: use tab-separated format
                row_text = " | ".join(row)
            lines.append(row_text)

        text = "\n".join(lines)
        print(f"    [{path.name}] -> {len(data_rows)} rows")
        return [("CSV Data", text)]
    except Exception as exc:
        print(f"    WARNING: could not read CSV {path.name} ({exc}) -- skipping")
        return []


def read_json_sections(path: Path) -> List[Tuple[str, str]]:
    """Flatten JSON structure to key: value lines."""
    import json

    def flatten_json(obj, prefix=""):
        """Recursively flatten JSON to list of 'path: value' lines."""
        lines = []
        if isinstance(obj, dict):
            for key, val in obj.items():
                new_prefix = f"{prefix}.{key}" if prefix else key
                if isinstance(val, (dict, list)):
                    lines.extend(flatten_json(val, new_prefix))
                else:
                    lines.append(f"{new_prefix}: {val}")
        elif isinstance(obj, list):
            for idx, item in enumerate(obj):
                new_prefix = f"{prefix}[{idx}]"
                if isinstance(item, (dict, list)):
                    lines.extend(flatten_json(item, new_prefix))
                else:
                    lines.append(f"{new_prefix}: {item}")
        return lines

    try:
        data = json.loads(path.read_text(encoding="utf-8"))
        lines = flatten_json(data)
        if lines:
            text = "\n".join(lines)
            print(f"    [{path.name}] -> {len(lines)} key-value pairs")
            return [("JSON Data", text)]
        return []
    except Exception as exc:
        print(f"    WARNING: could not parse JSON {path.name} ({exc}) -- skipping")
        return []


def read_pptx_sections(path: Path) -> List[Tuple[str, str]]:
    """Extract text from PowerPoint slides."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        sections: List[Tuple[str, str]] = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text_lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = normalize_text(shape.text)
                    if text:
                        slide_text_lines.append(text)

            if slide_text_lines:
                slide_text = "\n".join(slide_text_lines)
                sections.append((f"Slide {slide_idx}", slide_text))

        print(f"    [{path.name}] -> {len(sections)} slides with text")
        return sections
    except Exception as exc:
        print(f"    WARNING: could not read PPTX {path.name} ({exc}) -- skipping")
        return []


def read_excel_sections(path: Path) -> List[Tuple[str, str]]:
    """Read Excel file (.xlsx or .xls) — one section per sheet."""
    try:
        ext = path.suffix.lower()

        if ext == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path))
            sheets = wb.sheetnames
            sections: List[Tuple[str, str]] = []

            for sheet_name in sheets:
                ws = wb[sheet_name]
                rows_text = []
                for row in ws.iter_rows(values_only=True):
                    # Convert row tuple to text, handling None values
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_str.strip():
                        rows_text.append(row_str)

                if rows_text:
                    section_text = "\n".join(rows_text)
                    sections.append((f"Sheet: {sheet_name}", section_text))

            print(f"    [{path.name}] -> {len(sections)} sheet(s)")
            return sections

        elif ext == ".xls":
            import xlrd
            wb = xlrd.open_workbook(str(path))
            sections: List[Tuple[str, str]] = []

            for sheet_name in wb.sheet_names():
                ws = wb.sheet_by_name(sheet_name)
                rows_text = []
                for row_idx in range(ws.nrows):
                    row = ws.row_values(row_idx)
                    row_str = " | ".join(str(cell) if cell else "" for cell in row)
                    if row_str.strip():
                        rows_text.append(row_str)

                if rows_text:
                    section_text = "\n".join(rows_text)
                    sections.append((f"Sheet: {sheet_name}", section_text))

            print(f"    [{path.name}] -> {len(sections)} sheet(s)")
            return sections

        return []
    except Exception as exc:
        print(f"    WARNING: could not read Excel {path.name} ({exc}) -- skipping")
        return []


# ── Worker function for multiprocessing ───────────────────────────────────────
def process_file_for_chunks(path_str: str):
    """
    Worker: read a single file, split into sections, chunk text,
    and build payloads + texts for that file.

    Returns:
        (filename, n_sections, n_chunks, payloads, texts, elapsed_seconds)
    """
    path = Path(path_str)
    t0 = time.time()

    ext = path.suffix.lower()
    if ext == ".pdf":
        sections = read_pdf_sections(path)
    elif ext == ".docx":
        sections = read_docx_sections(path)
    elif ext == ".pptx":
        sections = read_pptx_sections(path)
    elif ext in (".xlsx", ".xls"):
        sections = read_excel_sections(path)
    elif ext == ".json":
        sections = read_json_sections(path)
    elif ext in (".txt", ".md", ".csv"):
        sections = read_csv_sections(path) if ext == ".csv" else read_txt_sections(path)
    else:
        print(f"    WARNING: unsupported file type {ext} -- skipping {path.name}")
        return path.name, 0, 0, [], [], 0

    doc_id = path.stem
    payloads = []
    texts = []
    file_chunks = 0

    for s_idx, (section_title, section_text) in enumerate(sections):
        chunks = chunk_text(section_text)
        file_chunks += len(chunks)
        for idx, chunk in enumerate(chunks):
            payloads.append({
                "doc_id": doc_id,
                "source_file": path.name,
                "section": section_title,
                "chunk_index": idx,
                "text": chunk,
            })
            texts.append(chunk)

    elapsed = time.time() - t0
    return path.name, len(sections), file_chunks, payloads, texts, elapsed


# ── Main ──────────────────────────────────────────────────────────────────────
def main() -> None:

    print("\n CUDA available:", torch.cuda.is_available())
    print("\n Device count:", torch.cuda.device_count())

    total_start = time.time()

    if not DATA_DIR.exists():
        raise SystemExit(f"Missing data folder: {DATA_DIR}")

    supported_formats = [".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".json", ".txt", ".md", ".csv"]
    files = [p for p in DATA_DIR.iterdir() if p.suffix.lower() in supported_formats]
    if not files:
        raise SystemExit(f"No supported files found in {DATA_DIR}. Supported: {', '.join(supported_formats)}")

    print(f"\n{'='*50}")
    print(f"  Found {len(files)} file(s) in {DATA_DIR}")
    for f in files:
        print(f"    * {f.name}  ({f.stat().st_size/1024:.1f} KB)")
    print(f"{'='*50}\n")

    # Step 1: Connect to Qdrant
    print(f"[1/4] Connecting to Qdrant at {QDRANT_HOST}:{QDRANT_PORT} ...")
    t0 = time.time()
    client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)
    collection_exists = client.collection_exists(COLLECTION)
    current_count = 0
    if collection_exists:
        current_count = client.get_collection(COLLECTION).points_count or 0
        print(f"      Collection '{COLLECTION}' exists with {current_count} points.")
    else:
        print(f"      Collection '{COLLECTION}' not found — will create once model is loaded.")
    print(f"      Connected in {time.time()-t0:.1f}s\n")

    # Step 2: Check hashes — skip unchanged files
    print(f"[2/4] Checking which files need ingestion ...")
    stored_hashes = load_hashes()

    # If collection is empty, force re-ingest regardless of hashes
    if current_count == 0 and stored_hashes:
        print("      Collection is empty but hashes exist — forcing full re-ingest.")
        stored_hashes = {}

    new_hashes = {}
    files_to_ingest: List[Path] = []

    for f in files:
        h = file_hash(f)
        new_hashes[f.name] = h
        if stored_hashes.get(f.name) == h:
            print(f"    SKIP  {f.name} (unchanged)")
        else:
            status = "NEW" if f.name not in stored_hashes else "CHANGED"
            print(f"    {status:7s} {f.name}")
            files_to_ingest.append(f)

    # Check for deleted files — remove their chunks from Qdrant
    deleted = [name for name in stored_hashes if name not in new_hashes]
    for name in deleted:
        print(f"    DELETED {name} -- removing from Qdrant...")
        client.delete(
            collection_name=COLLECTION,
            points_selector=Filter(
                must=[FieldCondition(key="source_file", match=MatchValue(value=name))]
            ),
        )

    if not files_to_ingest:
        print("\n  All files up to date -- nothing to ingest.")
        save_hashes(new_hashes)
        return

    print(f"\n  {len(files_to_ingest)} file(s) to ingest.\n")

    # Step 3: Load embedder (only when there is work to do)
    print(f"[3/4] Loading embedding model: {EMBED_MODEL_NAME} ...")
    t0 = time.time()
    from sentence_transformers import SentenceTransformer
    embedder = SentenceTransformer(EMBED_MODEL_NAME, device="cuda" if torch.cuda.is_available() else "cpu")
    dim = embedder.get_sentence_embedding_dimension()
    print(f"      Model loaded in {time.time()-t0:.1f}s  (dim={dim})\n")

    # Create collection now that we have the embedding dimension
    if not collection_exists:
        client.create_collection(
            collection_name=COLLECTION,
            vectors_config=VectorParams(size=dim, distance=Distance.COSINE),
        )
        print(f"      Collection '{COLLECTION}' created.\n")

    # Step 4: Parse, chunk & embed changed files
    print(f"[4/4] Parsing, chunking and embedding ...")

    # 4a) Remove old chunks for these files (single-threaded, cheap vs embedding)
    for f in files_to_ingest:
        print(f"  Deleting old chunks for {f.name} ...")
        client.delete(
            collection_name=COLLECTION,
            points_selector=Filter(
                must=[FieldCondition(key="source_file", match=MatchValue(value=f.name))]
            ),
        )

    # 4b) Process files in parallel (parsing + chunking)
    t0 = time.time()
    all_payloads = []
    all_texts = []

    max_workers = os.cpu_count() or 4
    print(f"\n  Using up to {max_workers} worker processes for parsing/chunking.\n")

    with ProcessPoolExecutor(max_workers=max_workers) as ex:
        futures = {
            ex.submit(process_file_for_chunks, str(f)): f.name
            for f in files_to_ingest
        }

        for fut in tqdm(as_completed(futures), total=len(futures), desc="  Processing files"):
            fname, n_sections, n_chunks, payloads, texts, elapsed = fut.result()
            print(f"  [{fname}] -> {n_sections} section(s), {n_chunks} chunk(s)  ({elapsed:.1f}s)")
            all_payloads.extend(payloads)
            all_texts.extend(texts)

    total_chars = sum(len(t) for t in all_texts)
    print(f"\n  Chunking done in {time.time()-t0:.1f}s")
    print(
        f"  Total: {len(all_texts)} chunks  |  ~{total_chars/1000:.0f}k chars  "
        f"|  avg {total_chars//max(len(all_texts),1)} chars/chunk\n"
    )

    # 4c) Embedding (GPU/CPU, but single process)
    print(f"  Embedding {len(all_texts)} chunks ...")
    t0 = time.time()
    vectors = embedder.encode(
        all_texts,
        batch_size=BATCH_SIZE,
        show_progress_bar=True,
        normalize_embeddings=True,
        device="cuda" if torch.cuda.is_available() else "cpu",
    )

    print(f"  Embedded in {time.time()-t0:.1f}s\n  Upserting ...")
    t0 = time.time()
    points = []
    batches_sent = 0

    for i in tqdm(range(len(all_payloads)), desc="  Upserting"):
        points.append(PointStruct(
            id=uuid.uuid4(),
            vector=vectors[i].tolist(),
            payload=all_payloads[i],
        ))
        if len(points) >= 256:
            client.upsert(collection_name=COLLECTION, points=points)
            batches_sent += 1
            points = []

    if points:
        client.upsert(collection_name=COLLECTION, points=points)
        batches_sent += 1

    print(f"  Upserted in {time.time()-t0:.1f}s  ({batches_sent} batch(es))\n")

    # Save updated hashes
    save_hashes(new_hashes)

    final_count = client.get_collection(COLLECTION).points_count or 0
    print(f"{'='*50}")
    print(f"  All done in {time.time()-total_start:.1f}s")
    print(f"  Collection : '{COLLECTION}'")
    print(f"  Points     : {final_count}")
    print(f"  Ingested   : {len(files_to_ingest)} file(s)")
    print(f"  Skipped    : {len(files) - len(files_to_ingest)} file(s)")
    print(f"{'='*50}\n")


if __name__ == "__main__":
    main()